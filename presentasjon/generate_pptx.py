#!/usr/bin/env python3
"""
Genererer Naturportrett_prototype.pptx — pedagogisk presentasjon for leder.

Iterasjon 2: rolig, elegant utforming med smal palett, hvite bokser med
fargede kantlinjer, generøs luft og logo med god plass.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ──────────────────────────────────────────────────────────────────────────────
# Oslo kommune designprofil — desaturert palett for ro og eleganse
# ──────────────────────────────────────────────────────────────────────────────

OSLO = {
    'morkebla': RGBColor(0x2A, 0x28, 0x59),     # primær - overskrifter, kant, mørke slides
    'lysebla_acc': RGBColor(0xB3, 0xF5, 0xFF),  # sparsomt brukt
    'morkegron': RGBColor(0x03, 0x4B, 0x45),    # "på plass / ✓"
    'lysgron': RGBColor(0xC7, 0xF6, 0xC9),      # subtil fyll for godkjent
    'gron_acc': RGBColor(0x43, 0xF8, 0xB6),     # aksentstriper, sparsomt
    'morkbeige': RGBColor(0xD0, 0xBF, 0xAE),    # myke avgrensninger
    'lysbeige': RGBColor(0xF8, 0xF0, 0xDD),     # subtil alternativ bakgrunn
    'gul': RGBColor(0xF9, 0xC6, 0x6B),          # kun advarsler
    'svart': RGBColor(0x2C, 0x2C, 0x2C),
    'hvit': RGBColor(0xFF, 0xFF, 0xFF),
    'gra': RGBColor(0xBB, 0xBB, 0xBB),
    'morkgra': RGBColor(0x55, 0x55, 0x55),
    'lysgra': RGBColor(0xEE, 0xEE, 0xEE),
}

FONT_PRIMARY = 'Oslo Sans'

ROOT = Path(__file__).resolve().parent.parent
LOGO_DARK = ROOT / 'Oslo visuell identitet' / 'Oslologo_skjerm_og_utskrift' / 'Oslo-logo-morkeblaa-til-skjerm-og-utskrift' / 'Oslo-logo-morkeblaa-RGB.png'
LOGO_WHITE = ROOT / 'Oslo visuell identitet' / 'Oslologo_skjerm_og_utskrift' / 'Oslo-logo-hvit-til-skjerm-og-utskrift' / 'Oslo-logo-hvit-RGB.png'

# Skjermbilder fra live-prototypen (lagres manuelt av brukeren)
SCREENSHOT_DIR = ROOT / 'presentasjon' / 'skjermbilder'
SCREENSHOT_ADDRESS = SCREENSHOT_DIR / '01_adresse.png'
SCREENSHOT_NATURPORTRETT = SCREENSHOT_DIR / '02_naturportrett.png'

OUTPUT = ROOT / 'presentasjon' / 'Naturportrett_prototype.pptx'

# ──────────────────────────────────────────────────────────────────────────────
# Layout-konstanter
# ──────────────────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Trygge marginer
MARGIN = Inches(0.7)
# Sone reservert for logo (nederst-venstre)
LOGO_W = Inches(1.3)
LOGO_H = Inches(0.65)
LOGO_X = MARGIN
LOGO_Y = SLIDE_H - LOGO_H - Inches(0.5)
# Innhold må stoppe minst denne y-verdien for å ikke kollidere med logo
CONTENT_Y_LIMIT = LOGO_Y - Inches(0.2)


# ──────────────────────────────────────────────────────────────────────────────
# Hjelpefunksjoner
# ──────────────────────────────────────────────────────────────────────────────

def make_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=None,
             font=FONT_PRIMARY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    return tb


def add_bulleted(slide, x, y, w, h, bullets, *, size=14, color=None,
                 line_spacing=1.35, after_pt=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    color = color or OSLO['svart']
    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(after_pt)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = '•  ' + item
        run.font.name = FONT_PRIMARY
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def outlined_box(slide, x, y, w, h, *, border_color=None, border_w=1.5,
                  fill=None, radius=0.06):
    """Standard mønster: hvit boks med farget tynn kant."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = radius
    if fill is None:
        box.fill.solid()
        box.fill.fore_color.rgb = OSLO['hvit']
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    border_color = border_color or OSLO['morkebla']
    box.line.color.rgb = border_color
    box.line.width = Pt(border_w)
    box.text_frame.margin_left = Pt(10)
    box.text_frame.margin_right = Pt(10)
    box.text_frame.margin_top = Pt(8)
    box.text_frame.margin_bottom = Pt(8)
    return box


def solid_box(slide, x, y, w, h, *, fill, border_color=None, radius=0.06):
    """Fylt boks (sjelden brukt — kun til hub/aksenter)."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = radius
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if border_color is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = border_color
        box.line.width = Pt(1.0)
    box.text_frame.margin_left = Pt(10)
    box.text_frame.margin_right = Pt(10)
    box.text_frame.margin_top = Pt(8)
    box.text_frame.margin_bottom = Pt(8)
    return box


def fill_box_text(box, lines, *, size=14, bold=False, italic=False,
                  color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                  font=FONT_PRIMARY, line_spacing=1.25):
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    color = color or OSLO['svart']
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ''
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(line, dict):
            run = p.add_run()
            run.text = line['text']
            run.font.name = font
            run.font.size = Pt(line.get('size', size))
            run.font.bold = line.get('bold', bold)
            run.font.italic = line.get('italic', italic)
            run.font.color.rgb = line.get('color', color)
        else:
            run = p.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color


def accent_stripe(slide, x, y, w, *, color=None, thickness=Inches(0.07)):
    """Tynn aksentstripe øverst i et kort."""
    color = color or OSLO['gron_acc']
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thickness)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    return stripe


def add_arrow(slide, x1, y1, x2, y2, *, color=None, width=1.5,
              dashed=False):
    color = color or OSLO['morkebla']
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line_xml = line.line._get_or_add_ln()
    tail = etree.SubElement(line_xml, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    if dashed:
        # legg til prstDash dash for stiplet
        dash = etree.SubElement(line_xml, qn('a:prstDash'))
        dash.set('val', 'dash')
    return line


def add_logo(slide, *, white=False):
    """Logo deaktivert globalt — vises ikke i denne versjonen."""
    return None


def set_slide_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def title_header(slide, title, subtitle=None):
    add_text(slide, MARGIN, Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, size=30, bold=True, color=OSLO['morkebla'])
    if subtitle:
        add_text(slide, MARGIN, Inches(1.2),
                 SLIDE_W - 2 * MARGIN, Inches(0.45),
                 subtitle, size=16, color=OSLO['morkgra'])
    # Tynn linje under
    line = slide.shapes.add_connector(1,
                                       MARGIN, Inches(1.75),
                                       SLIDE_W - MARGIN, Inches(1.75))
    line.line.color.rgb = OSLO['morkbeige']
    line.line.width = Pt(0.75)


def page_number(slide, n, total, *, dark=False):
    color = OSLO['hvit'] if dark else OSLO['morkgra']
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.5),
             Inches(0.7), Inches(0.3),
             f'{n} / {total}', size=10, color=color, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# Slides
# ──────────────────────────────────────────────────────────────────────────────

BLANK = None
TOTAL = 16


def s01_title(prs):
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, OSLO['morkebla'])

    add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.5),
             'Naturportrett', size=72, bold=True, color=OSLO['hvit'])
    add_text(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.6),
             'KI-drevet vurdering av biologisk mangfold i Oslo',
             size=22, color=OSLO['lysebla_acc'])

    accent_stripe(s, Inches(0.9), Inches(2.92), Inches(2.5),
                   color=OSLO['gron_acc'])

    add_text(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.4),
             'Prototype • Plan- og bygningsetaten • Mai 2026',
             size=15, color=OSLO['lysebla_acc'])
    add_text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.4),
             'Andreas Fadum Haugstad',
             size=15, color=OSLO['hvit'], bold=True)

    add_logo(s, white=True)


def s02_why(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Hvorfor Naturportrett?')

    items = [
        ('Tidlig naturhensyn',
         'Plan- og byggesaker krever vurdering av biologisk mangfold tidlig — før det blir kostbart å endre kurs.'),
        ('Fragmentert kunnskap',
         'Informasjon ligger spredt på Artsdatabanken, iNaturalist, GBIF, kommunale rapporter og lovverk. Tidkrevende å samle.'),
        ('KI sammenstiller raskt',
         'KI kan på sekunder lese flere kilder og levere et strukturert utgangspunkt som fagperson kvalitetssikrer.'),
    ]
    card_w = Inches(3.7)
    card_h = Inches(3.2)
    gap = Inches(0.5)
    total = 3 * card_w + 2 * gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.4)
    for i, (title, body) in enumerate(items):
        x = start_x + i * (card_w + gap)
        outlined_box(s, x, y, card_w, card_h, border_color=OSLO['morkebla'],
                     border_w=1.5)
        # Nummer
        add_text(s, x + Inches(0.4), y + Inches(0.4),
                 Inches(0.6), Inches(0.6),
                 f'0{i+1}', size=22, bold=True, color=OSLO['morkebla'])
        add_text(s, x + Inches(0.4), y + Inches(1.05),
                 card_w - Inches(0.8), Inches(0.6),
                 title, size=18, bold=True, color=OSLO['morkebla'])
        add_text(s, x + Inches(0.4), y + Inches(1.75),
                 card_w - Inches(0.8), Inches(1.4),
                 body, size=13, color=OSLO['svart'], line_spacing=1.4)

    add_logo(s)
    page_number(s, 2, TOTAL)


def s03_products(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'De fire produktene',
                 subtitle='Et naturportrett er hovedproduktet — de tre andre er detaljutdypninger')

    products = [
        ('🌐', 'Naturportrett',
         'Områdeoversikt: kart, viktige arter, naturtyper, økologiske sammenhenger, trusler.',
         OSLO['morkebla']),
        ('🌳', 'Naturtypeportrett',
         'Detaljert beskrivelse av en konkret naturtype (NiN-kode), inkl. typiske arter og hensyn.',
         OSLO['morkegron']),
        ('🦜', 'Artsportrett',
         'Detaljert dyreportrett: levevis, årssyklus, næringskilder, designtiltak.',
         OSLO['morkebla']),
        ('🌿', 'Planteportrett',
         'Detaljert planteportrett: habitatkrav, pollinator-verdi, samplanting, erfaring i Norge.',
         OSLO['morkegron']),
    ]
    card_w = Inches(2.7)
    card_h = Inches(3.4)
    gap = Inches(0.35)
    total = 4 * card_w + 3 * gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.2)
    for i, (icon, title, body, stripe_color) in enumerate(products):
        x = start_x + i * (card_w + gap)
        outlined_box(s, x, y, card_w, card_h,
                     border_color=OSLO['morkebla'], border_w=1.5)

        add_text(s, x + Inches(0.3), y + Inches(0.45),
                 card_w - Inches(0.6), Inches(0.9),
                 icon, size=38, color=OSLO['morkebla'])
        add_text(s, x + Inches(0.3), y + Inches(1.5),
                 card_w - Inches(0.6), Inches(0.5),
                 title, size=16, bold=True, color=OSLO['morkebla'])
        add_text(s, x + Inches(0.3), y + Inches(2.05),
                 card_w - Inches(0.6), Inches(1.2),
                 body, size=12, color=OSLO['svart'], line_spacing=1.4)

    add_text(s, MARGIN, Inches(5.9), SLIDE_W - 2 * MARGIN, Inches(0.4),
             'Et naturportrett kan inneholde flere naturtyper, arter og planter — som hver kan få sitt eget detaljerte portrett.',
             size=13, italic=True, color=OSLO['morkgra'],
             align=PP_ALIGN.CENTER)

    add_logo(s)
    page_number(s, 3, TOTAL)


def s04_portrait_content(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Hva inneholder et portrett?',
                 subtitle='Eksempel: oppbygning av et artsportrett')

    mock_x = Inches(1.5)
    mock_y = Inches(2.1)
    mock_w = Inches(10.3)
    mock_h = Inches(4.4)
    outlined_box(s, mock_x, mock_y, mock_w, mock_h,
                  border_color=OSLO['morkebla'], border_w=1.5)
    # Mockup-overskrift
    add_text(s, mock_x + Inches(0.4), mock_y + Inches(0.25),
             Inches(4), Inches(0.5),
             'Artsportrett', size=20, bold=True, color=OSLO['morkebla'])
    add_text(s, mock_x + mock_w - Inches(2.0), mock_y + Inches(0.3),
             Inches(1.7), Inches(0.4),
             '17. mai 2026', size=10, color=OSLO['morkgra'],
             align=PP_ALIGN.RIGHT)

    # Navn-boks + vit-navn-boks
    outlined_box(s, mock_x + Inches(0.4), mock_y + Inches(0.95),
                  Inches(3.2), Inches(0.7),
                  border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(0.55), mock_y + Inches(1.07),
             Inches(3.0), Inches(0.45),
             'Folkenavn', size=13, bold=True, color=OSLO['svart'])

    outlined_box(s, mock_x + Inches(0.4), mock_y + Inches(1.8),
                  Inches(3.2), Inches(0.7),
                  border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(0.55), mock_y + Inches(1.92),
             Inches(3.0), Inches(0.45),
             'Vitenskapelig navn', size=13, italic=True, color=OSLO['svart'])

    # Status
    outlined_box(s, mock_x + Inches(3.85), mock_y + Inches(0.95),
                  Inches(1.9), Inches(1.55),
                  border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(4.0), mock_y + Inches(1.05),
             Inches(1.7), Inches(0.4),
             'Rødliste status', size=11, bold=True, color=OSLO['svart'])
    add_text(s, mock_x + Inches(4.0), mock_y + Inches(1.7),
             Inches(1.7), Inches(0.6),
             'VU', size=22, bold=True, color=OSLO['gul'])

    # Foto-plassholder
    photo = outlined_box(s, mock_x + Inches(6.0), mock_y + Inches(0.95),
                         Inches(3.9), Inches(2.4),
                         border_color=OSLO['svart'], border_w=1.0,
                         fill=OSLO['lysbeige'])
    fill_box_text(photo, '[Foto av arten]', size=13, italic=True,
                  color=OSLO['morkgra'], align=PP_ALIGN.CENTER)

    # Karakteristikker
    outlined_box(s, mock_x + Inches(0.4), mock_y + Inches(2.65),
                  Inches(5.35), Inches(0.7),
                  border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(0.55), mock_y + Inches(2.77),
             Inches(5.0), Inches(0.45),
             'Karakteristikker — artsfamilie · utbredelse',
             size=12, color=OSLO['svart'])

    # Beskrivelse + habitat
    outlined_box(s, mock_x + Inches(0.4), mock_y + Inches(3.5),
                  Inches(4.55), Inches(0.7),
                  border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(0.55), mock_y + Inches(3.62),
             Inches(4.3), Inches(0.45),
             'Beskrivelse — størrelse, farger, kjønnsforskjeller',
             size=12, color=OSLO['svart'])

    outlined_box(s, mock_x + Inches(5.15), mock_y + Inches(3.5),
                  Inches(4.75), Inches(0.7),
                  border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(5.3), mock_y + Inches(3.62),
             Inches(4.45), Inches(0.45),
             'Foretrukne habitater', size=12, color=OSLO['svart'])

    add_text(s, MARGIN, Inches(6.7), SLIDE_W - 2 * MARGIN, Inches(0.3),
             'Portrettet har også årssyklus, næringskilder, attributter (nøkkelart/ansvarsart) og praktiske designtiltak.',
             size=12, italic=True, color=OSLO['morkgra'],
             align=PP_ALIGN.CENTER)

    add_logo(s)
    page_number(s, 4, TOTAL)


def s05_user_flow(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Brukerflyt — fem steg',
                 subtitle='Fra adresse til ferdig portrett klart for nedlasting som PDF')

    steps = ['Adresse', 'Naturportrett', 'Velg portrettype',
             'Velg subjekt', 'Detaljportrett']
    box_w = Inches(2.05)
    box_h = Inches(2.2)
    gap = Inches(0.25)
    total = 5 * box_w + 4 * gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.6)

    for i, label in enumerate(steps):
        x = start_x + i * (box_w + gap)
        # Siste trinn får subtil grønn fyll for å indikere "mål"
        fill = OSLO['lysgron'] if i == 4 else None
        outlined_box(s, x, y, box_w, box_h,
                     border_color=OSLO['morkebla'], border_w=1.5,
                     fill=fill)
        # Nummer-sirkel
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + box_w / 2 - Inches(0.4),
                                     y + Inches(0.3),
                                     Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = OSLO['morkebla']
        circle.line.fill.background()
        fill_box_text(circle, str(i + 1), size=24, bold=True,
                      color=OSLO['hvit'], align=PP_ALIGN.CENTER)

        add_text(s, x + Inches(0.1), y + Inches(1.35),
                 box_w - Inches(0.2), Inches(0.6),
                 label, size=14, bold=True, color=OSLO['morkebla'],
                 align=PP_ALIGN.CENTER)

        if i < 4:
            ax = x + box_w
            ay = y + box_h / 2
            add_arrow(s, ax + Inches(0.03), ay,
                      ax + gap - Inches(0.03), ay,
                      color=OSLO['morkebla'], width=2.0)

    add_text(s, MARGIN, Inches(5.6), SLIDE_W - 2 * MARGIN, Inches(0.5),
             '↻  Brukeren kan generere flere detaljportretter for samme adresse',
             size=13, italic=True, color=OSLO['morkgra'],
             align=PP_ALIGN.CENTER)

    add_logo(s)
    page_number(s, 5, TOTAL)


def add_screenshot(slide, path, *, fallback_caption):
    """Embed skjermbilde sentrert under tittelen, eller vis pen plassholder."""
    # Område tilgjengelig for skjermbildet
    avail_x = Inches(0.7)
    avail_y = Inches(2.0)
    avail_w = SLIDE_W - 2 * Inches(0.7)
    avail_h = Inches(5.0)

    if not path.exists():
        # Plassholder
        outlined_box(slide, avail_x, avail_y, avail_w, avail_h,
                     border_color=OSLO['morkbeige'], border_w=1.0,
                     fill=OSLO['lysbeige'])
        add_text(slide, avail_x + Inches(0.5),
                 avail_y + avail_h / 2 - Inches(0.3),
                 avail_w - Inches(1.0), Inches(0.6),
                 fallback_caption, size=14, italic=True,
                 color=OSLO['morkgra'], align=PP_ALIGN.CENTER)
        return

    # Få faktiske dimensjoner med Pillow
    try:
        from PIL import Image
        with Image.open(path) as img:
            iw, ih = img.size
    except Exception:
        iw, ih = 1600, 2400  # antagelse hvis Pillow feiler

    # Skaler så det passer i avail_w x avail_h, behold sideforhold
    avail_w_emu = int(avail_w)
    avail_h_emu = int(avail_h)
    img_ratio = iw / ih
    avail_ratio = avail_w_emu / avail_h_emu
    if img_ratio > avail_ratio:
        # bildet er bredere — fyll bredden
        new_w = avail_w_emu
        new_h = int(avail_w_emu / img_ratio)
    else:
        # bildet er høyere — fyll høyden
        new_h = avail_h_emu
        new_w = int(avail_h_emu * img_ratio)

    # Sentrer
    x = avail_x + (avail_w_emu - new_w) // 2
    y = avail_y + (avail_h_emu - new_h) // 2

    pic = slide.shapes.add_picture(str(path), x, y, width=new_w, height=new_h)
    # Tynn ramme rundt bildet
    pic.line.color.rgb = OSLO['morkbeige']
    pic.line.width = Pt(0.75)


def s06_demo_address(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Slik fungerer adressesøket',
                 subtitle='Brukeren søker etter en adresse i Oslo og får automatiske forslag fra Kartverket')

    add_screenshot(s, SCREENSHOT_ADDRESS,
                   fallback_caption='Lagre skjermbilde her: presentasjon/skjermbilder/01_adresse.png')

    page_number(s, 6, TOTAL)


def s07_demo_naturportrett(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Slik ser naturportrettet ut',
                 subtitle='Områdeoversikt med kart, naturtyper, arter og tekstlige seksjoner — Vahls gate 1, Oslo')

    add_screenshot(s, SCREENSHOT_NATURPORTRETT,
                   fallback_caption='Lagre skjermbilde her: presentasjon/skjermbilder/02_naturportrett.png')

    page_number(s, 7, TOTAL)


def _s07_demo_naturportrett_OLD(prs):
    """Beholdt som referanse — ikke i bruk."""
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Demo — Steg 2: Naturportrettet genereres automatisk',
                 subtitle='Vahls gate 1, 0187 Oslo • innenfor 500 m radius')

    mock_x = Inches(0.7)
    mock_y = Inches(2.0)
    mock_w = SLIDE_W - 2 * Inches(0.7)
    mock_h = Inches(4.7)
    outlined_box(s, mock_x, mock_y, mock_w, mock_h,
                 border_color=OSLO['morkbeige'], border_w=1.0,
                 fill=OSLO['hvit'])

    # Søker innenfor-banner (lysegrønn)
    banner_y = mock_y + Inches(0.15)
    outlined_box(s, mock_x + Inches(0.2), banner_y,
                 mock_w - Inches(0.4), Inches(0.55),
                 border_color=OSLO['morkegron'], border_w=1.0,
                 fill=OSLO['lysgron'])
    add_text(s, mock_x + Inches(0.4), banner_y + Inches(0.12),
             Inches(11), Inches(0.35),
             '📍   Søker innenfor 500 m fra Vahls gate 1, 0187 OSLO',
             size=12, bold=True, color=OSLO['morkegron'])

    # Venstre: kart-mockup
    map_x = mock_x + Inches(0.2)
    map_y = mock_y + Inches(0.95)
    map_w = Inches(5.5)
    map_h = Inches(3.5)
    outlined_box(s, map_x, map_y, map_w, map_h,
                 border_color=OSLO['svart'], border_w=1.0,
                 fill=OSLO['lysbeige'])
    # Kartsirkel
    cx = map_x + map_w / 2
    cy = map_y + map_h / 2
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - Inches(1.4), cy - Inches(1.4),
                                 Inches(2.8), Inches(2.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = OSLO['lysgron']
    circle.line.color.rgb = OSLO['morkebla']
    circle.line.width = Pt(2)
    # Marker
    marker = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - Inches(0.12), cy - Inches(0.12),
                                 Inches(0.24), Inches(0.24))
    marker.fill.solid()
    marker.fill.fore_color.rgb = OSLO['morkebla']
    marker.line.color.rgb = OSLO['hvit']
    marker.line.width = Pt(1.5)
    add_text(s, map_x + Inches(0.15), map_y + Inches(0.1),
             Inches(3), Inches(0.3),
             'Oversiktskart', size=11, bold=True, color=OSLO['morkebla'])
    add_text(s, map_x + Inches(0.15), map_y + map_h - Inches(0.35),
             map_w - Inches(0.3), Inches(0.3),
             '500 m influenssone (markør = adresse)',
             size=9, italic=True, color=OSLO['morkgra'])

    # Høyre: arts-grid
    grid_x = mock_x + Inches(5.95)
    grid_y = mock_y + Inches(0.95)
    grid_w = mock_w - Inches(6.15)
    add_text(s, grid_x, grid_y,
             grid_w, Inches(0.35),
             'Natur og arter i nærområdet', size=13, bold=True,
             color=OSLO['morkebla'])

    species = [
        ('Spisslønn', 'Acer platanoides', 'PLANTE', 'SE', OSLO['svart']),
        ('Ringdue', 'Columba palumbus', 'FUGL', 'LC', OSLO['morkegron']),
        ('Harlekinmarihøne', 'Harmonia axyridis', 'INSEKT', 'SE', OSLO['svart']),
        ('Løkurt', 'Alliaria petiolata', 'PLANTE', 'SE', OSLO['svart']),
        ('Stillits', 'Carduelis carduelis', 'FUGL', 'LC', OSLO['morkegron']),
        ('Bjørk', 'Betula pendula', 'PLANTE', 'LC', OSLO['morkegron']),
    ]
    # 2-kolonner kort-grid (kompakt)
    cols = 2
    card_w = (grid_w - Inches(0.15)) / 2
    card_h = Inches(0.95)
    row_gap = Inches(0.1)
    for i, (no, sci, cat, status, status_color) in enumerate(species):
        c = i % cols
        r = i // cols
        x = grid_x + c * (card_w + Inches(0.15))
        y = grid_y + Inches(0.45) + r * (card_h + row_gap)
        outlined_box(s, x, y, card_w, card_h,
                     border_color=OSLO['morkbeige'], border_w=1.0)
        # Foto-plassholder
        photo = outlined_box(s, x + Inches(0.08), y + Inches(0.1),
                             Inches(0.75), Inches(0.75),
                             border_color=OSLO['lysgra'], border_w=0.5,
                             fill=OSLO['lysbeige'])
        fill_box_text(photo, '🖼', size=14, color=OSLO['morkgra'],
                      align=PP_ALIGN.CENTER)
        # Tekst
        text_x = x + Inches(0.95)
        text_w = card_w - Inches(1.05)
        add_text(s, text_x, y + Inches(0.08),
                 text_w, Inches(0.3),
                 no, size=12, bold=True, color=OSLO['svart'])
        add_text(s, text_x, y + Inches(0.35),
                 text_w, Inches(0.25),
                 sci, size=9, italic=True, color=OSLO['morkgra'])
        # Badges
        add_text(s, text_x, y + Inches(0.62),
                 Inches(0.8), Inches(0.25),
                 cat, size=9, bold=True, color=OSLO['morkebla'])
        add_text(s, text_x + Inches(0.85), y + Inches(0.62),
                 Inches(0.6), Inches(0.25),
                 status, size=10, bold=True, color=status_color)

    add_logo(s)
    page_number(s, 7, TOTAL)


def s08_demo_detail(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Demo — Steg 3: Detaljportrett + PDF',
                 subtitle='Eksempel: Artsportrett for Ringdue (Columba palumbus)')

    mock_x = Inches(0.7)
    mock_y = Inches(2.0)
    mock_w = SLIDE_W - 2 * Inches(0.7)
    mock_h = Inches(4.6)
    outlined_box(s, mock_x, mock_y, mock_w, mock_h,
                 border_color=OSLO['morkbeige'], border_w=1.0,
                 fill=OSLO['hvit'])

    # Tittel-rad
    add_text(s, mock_x + Inches(0.3), mock_y + Inches(0.2),
             Inches(5), Inches(0.5),
             'Artsportrett', size=24, color=OSLO['morkebla'])
    add_text(s, mock_x + mock_w - Inches(2.0), mock_y + Inches(0.25),
             Inches(1.7), Inches(0.35),
             '17. mai 2026', size=10, color=OSLO['morkgra'],
             align=PP_ALIGN.RIGHT)

    # Hero-rad: navn + status + foto
    hero_y = mock_y + Inches(0.95)
    # Navn-boks
    outlined_box(s, mock_x + Inches(0.3), hero_y,
                 Inches(3.2), Inches(0.65),
                 border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(0.45), hero_y + Inches(0.13),
             Inches(3.0), Inches(0.4),
             'Ringdue', size=15, bold=True, color=OSLO['svart'])
    # Vit-navn
    outlined_box(s, mock_x + Inches(0.3), hero_y + Inches(0.8),
                 Inches(3.2), Inches(0.65),
                 border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(0.45), hero_y + Inches(0.93),
             Inches(3.0), Inches(0.4),
             'Columba palumbus', size=13, italic=True, color=OSLO['svart'])
    # Status
    outlined_box(s, mock_x + Inches(3.65), hero_y,
                 Inches(1.5), Inches(1.45),
                 border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(3.78), hero_y + Inches(0.1),
             Inches(1.3), Inches(0.35),
             'Rødliste status', size=10, bold=True, color=OSLO['svart'])
    add_text(s, mock_x + Inches(3.78), hero_y + Inches(0.7),
             Inches(1.3), Inches(0.5),
             'LC', size=22, bold=True, color=OSLO['morkegron'])
    # Foto
    photo = outlined_box(s, mock_x + Inches(5.3), hero_y,
                         Inches(3.0), Inches(1.45),
                         border_color=OSLO['svart'], border_w=1.0,
                         fill=OSLO['lysbeige'])
    fill_box_text(photo, '[Foto: Ringdue]', size=11, italic=True,
                  color=OSLO['morkgra'], align=PP_ALIGN.CENTER)

    # Detalj-bokser (matcher portrait-doc__factbox-inline)
    detail_y = mock_y + Inches(2.55)
    details = [
        ('Karakteristikker', 'Columbidae · hele Norge'),
        ('Foretrukne habitater', 'Skoger, parker, kulturlandskap, urbane områder'),
        ('Næringskilder', 'Frø, knopper, bær, korn'),
    ]
    for i, (head, body) in enumerate(details):
        y = detail_y + i * Inches(0.65)
        outlined_box(s, mock_x + Inches(0.3), y,
                     mock_w - Inches(0.6), Inches(0.5),
                     border_color=OSLO['svart'], border_w=1.0)
        add_text(s, mock_x + Inches(0.5), y + Inches(0.1),
                 Inches(3.5), Inches(0.3),
                 head, size=11, bold=True, color=OSLO['svart'])
        add_text(s, mock_x + Inches(4.0), y + Inches(0.1),
                 mock_w - Inches(4.3), Inches(0.3),
                 body, size=11, color=OSLO['svart'])

    # Designtiltak-blokk
    tiltak_y = detail_y + 3 * Inches(0.65) + Inches(0.05)
    outlined_box(s, mock_x + Inches(0.3), tiltak_y,
                 mock_w - Inches(0.6), Inches(0.7),
                 border_color=OSLO['svart'], border_w=1.0)
    add_text(s, mock_x + Inches(0.5), tiltak_y + Inches(0.1),
             Inches(3.5), Inches(0.3),
             'Praktiske designtiltak',
             size=11, bold=True, color=OSLO['svart'])
    add_text(s, mock_x + Inches(0.5), tiltak_y + Inches(0.38),
             mock_w - Inches(0.8), Inches(0.3),
             '• bevare høye trær  • skjul-vegetasjon  • unngå glass mot trekkretningen',
             size=11, color=OSLO['svart'])

    # PDF-knapp
    btn_w = Inches(2.6)
    btn_y = mock_y + mock_h - Inches(0.6)
    btn_x = mock_x + mock_w - btn_w - Inches(0.3)
    solid_box(s, btn_x, btn_y, btn_w, Inches(0.45),
              fill=OSLO['morkegron'], radius=0.4)
    add_text(s, btn_x, btn_y + Inches(0.08),
             btn_w, Inches(0.3),
             '📄  Last ned som PDF', size=12, bold=True,
             color=OSLO['hvit'], align=PP_ALIGN.CENTER)

    add_logo(s)
    page_number(s, 8, TOTAL)


def s09_architecture(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Bak kulissene — tre lag',
                 subtitle='Hva som skjer fra knappen blir trykket til portrettet vises')

    cols = [
        ('Nettleser', 'Det brukeren ser:\nadresse-søk, kart, portretter, knapper'),
        ('Server', 'Henter data fra åpne kilder, snakker med KI, formaterer resultatet'),
        ('KI', 'Skriver portretttekst basert på dataene som leveres'),
    ]
    col_w = Inches(3.2)
    col_h = Inches(3.4)
    gap = Inches(0.7)
    total = 3 * col_w + 2 * gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.4)
    for i, (head, body) in enumerate(cols):
        x = start_x + i * (col_w + gap)
        outlined_box(s, x, y, col_w, col_h,
                     border_color=OSLO['morkebla'], border_w=1.5)
        # Sirkel-nummer
        num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + col_w / 2 - Inches(0.45),
                                  y + Inches(0.45),
                                  Inches(0.9), Inches(0.9))
        num.fill.solid()
        num.fill.fore_color.rgb = OSLO['morkebla']
        num.line.fill.background()
        fill_box_text(num, str(i + 1), size=28, bold=True,
                      color=OSLO['hvit'], align=PP_ALIGN.CENTER)

        add_text(s, x + Inches(0.3), y + Inches(1.55),
                 col_w - Inches(0.6), Inches(0.5),
                 head, size=18, bold=True, color=OSLO['morkebla'],
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(2.15),
                 col_w - Inches(0.6), Inches(1.1),
                 body, size=13, color=OSLO['svart'],
                 align=PP_ALIGN.CENTER, line_spacing=1.4)

        # Pil mellom kolonner
        if i < 2:
            ax = x + col_w
            ay = y + col_h / 2
            add_arrow(s, ax + Inches(0.08), ay,
                      ax + gap - Inches(0.08), ay,
                      color=OSLO['morkebla'], width=2.0)

    add_text(s, MARGIN, Inches(6.2), SLIDE_W - 2 * MARGIN, Inches(0.4),
             'Hemmelig informasjon (som tilgangsnøkler) ligger kun på serveren — aldri synlig i nettleseren.',
             size=12, italic=True, color=OSLO['morkgra'],
             align=PP_ALIGN.CENTER)

    page_number(s, 8, TOTAL)


def s10_data_sources(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Datakilder i prototypen',
                 subtitle='Hvor informasjonen kommer fra — og hva som er gratis vs. betalt')

    # Hub
    hub_w = Inches(2.6)
    hub_h = Inches(1.0)
    hub_x = SLIDE_W / 2 - hub_w / 2
    hub_y = Inches(4.3)
    hub = solid_box(s, hub_x, hub_y, hub_w, hub_h,
                     fill=OSLO['morkebla'])
    fill_box_text(hub, 'Naturportrett',
                  size=18, bold=True, color=OSLO['hvit'],
                  align=PP_ALIGN.CENTER)

    # Datakilder
    sources = [
        ('Kartverket', 'adresse, koordinater', 'gratis',
         Inches(0.5), Inches(2.3)),
        ('iNaturalist + GBIF', 'arter, foto', 'gratis',
         Inches(5.4), Inches(2.3)),
        ('Artsdatabanken', 'rødliste, fremmedarts', 'gratis',
         Inches(10.3), Inches(2.3)),
        ('OpenStreetMap', 'bakgrunnskart', 'gratis',
         Inches(0.5), Inches(5.8)),
        ('KI', 'sammenstiller informasjonen', 'betalt',
         Inches(10.3), Inches(5.8)),
    ]
    src_w = Inches(2.5)
    src_h = Inches(1.1)
    hub_center_x = hub_x + hub_w / 2
    hub_center_y = hub_y + hub_h / 2

    for label, what, pricing, x, y in sources:
        # Boks — subtil bakgrunn for KI-boksen så tekst er mer lesbar
        is_paid = pricing == 'betalt'
        border = OSLO['gul'] if is_paid else OSLO['morkegron']
        fill = OSLO['lysbeige'] if is_paid else None
        outlined_box(s, x, y, src_w, src_h,
                     border_color=border, border_w=1.5, fill=fill)
        add_text(s, x + Inches(0.15), y + Inches(0.1),
                 src_w - Inches(0.3), Inches(0.35),
                 label, size=12, bold=True, color=OSLO['morkebla'],
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(0.42),
                 src_w - Inches(0.3), Inches(0.3),
                 what, size=10, color=OSLO['svart'],
                 align=PP_ALIGN.CENTER)
        tag_color = OSLO['svart'] if is_paid else OSLO['morkegron']
        add_text(s, x + Inches(0.15), y + Inches(0.72),
                 src_w - Inches(0.3), Inches(0.3),
                 f'({pricing})', size=10, bold=True, italic=True,
                 color=tag_color, align=PP_ALIGN.CENTER)

        # Pil: start fra boksen sin nærmeste kant, ikke senter
        src_center_x = x + src_w / 2
        src_center_y = y + src_h / 2
        # Bestem startpunkt basert på hvilken side som er nærmest hub
        dx = hub_center_x - src_center_x
        dy = hub_center_y - src_center_y
        # Velg startkant: hvis hub er mest til høyre/venstre, start på horisontal kant
        if abs(dx) > abs(dy):
            if dx > 0:  # hub er til høyre
                start_x = x + src_w
                start_y = src_center_y
            else:
                start_x = x
                start_y = src_center_y
        else:
            if dy > 0:  # hub er nedenfor
                start_x = src_center_x
                start_y = y + src_h
            else:
                start_x = src_center_x
                start_y = y

        # Endepunkt: nærmeste kant på hub
        if abs(dx) > abs(dy):
            if dx > 0:
                end_x = hub_x
                end_y = hub_center_y
            else:
                end_x = hub_x + hub_w
                end_y = hub_center_y
        else:
            if dy > 0:
                end_x = hub_center_x
                end_y = hub_y
            else:
                end_x = hub_center_x
                end_y = hub_y + hub_h

        add_arrow(s, start_x, start_y, end_x, end_y,
                  color=OSLO['morkebla'], width=1.8)

    # Forklaring
    add_text(s, MARGIN, Inches(7.0), SLIDE_W - 2 * MARGIN, Inches(0.35),
             'Alle datakilder unntatt KI er gratis offentlige tjenester. KI er den eneste betalte tjenesten — koster < 0,15 USD per komplett sesjon.',
             size=11, italic=True, color=OSLO['morkgra'],
             align=PP_ALIGN.CENTER)

    page_number(s, 9, TOTAL)


def s11_ai_role(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Hva gjør KI — og hva gjør den ikke?',
                 subtitle='Det er viktig å forstå hva systemet kan og ikke kan')

    col_w = Inches(5.4)
    col_h = Inches(4.0)
    gap = Inches(0.5)
    total = 2 * col_w + gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.3)

    # Gjør
    outlined_box(s, start_x, y, col_w, col_h,
                 border_color=OSLO['morkegron'], border_w=2.0)
    add_text(s, start_x + Inches(0.4), y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.5),
             '✓  GJØR', size=22, bold=True, color=OSLO['morkegron'])
    add_bulleted(s, start_x + Inches(0.4), y + Inches(1.0),
                 col_w - Inches(0.8), col_h - Inches(1.2),
                 [
                     'Sammenstiller informasjon fra dataene den får',
                     'Skriver strukturert tekst med faste seksjoner',
                     'Følger Oslo kommunes naturmangfoldsstrategi',
                     'Foreslår praktiske designtiltak',
                     'Holder konsistent struktur uansett område',
                 ], size=13, color=OSLO['morkegron'], after_pt=10)

    # Gjør IKKE
    x2 = start_x + col_w + gap
    outlined_box(s, x2, y, col_w, col_h,
                 border_color=OSLO['morkbeige'], border_w=2.0)
    add_text(s, x2 + Inches(0.4), y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.5),
             '✗  GJØR IKKE', size=22, bold=True, color=OSLO['morkbeige'])
    add_bulleted(s, x2 + Inches(0.4), y + Inches(1.0),
                 col_w - Inches(0.8), col_h - Inches(1.2),
                 [
                     'Søker ikke selv etter ny informasjon på nett',
                     'Verifiserer ikke fakta mot offisielle databaser',
                     'Kan av og til hallusinere detaljer',
                     'Erstatter ikke fagperson — er et utgangspunkt',
                     'Tar ikke høyde for siste års endringer',
                 ], size=13, color=OSLO['svart'], after_pt=10)

    # Advarselsstripe i beige
    outlined_box(s, MARGIN, Inches(6.5),
                 SLIDE_W - 2 * MARGIN, Inches(0.55),
                 border_color=OSLO['morkbeige'], border_w=1.0,
                 fill=OSLO['lysbeige'])
    add_text(s, MARGIN, Inches(6.6), SLIDE_W - 2 * MARGIN, Inches(0.4),
             'Alle portretter må kvalitetssikres av fagperson før bruk i saksbehandling',
             size=14, bold=True, color=OSLO['morkebla'],
             align=PP_ALIGN.CENTER)

    page_number(s, 10, TOTAL)


def s12_more_sources(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Veikart — flere datakilder',
                 subtitle='Slik kan prototypen vokse. Hver ny kilde gir bedre presisjon.')

    existing = [
        'Kartverket', 'iNaturalist', 'GBIF',
        'OpenStreetMap', 'Artsdatabanken (statisk)',
    ]
    future = [
        'Artsdatabanken API (live)',
        'NiN-systemet (naturtyper)',
        'Norge i bilder (ortofoto)',
        'Miljødirektoratet (vern, vilt)',
        'Plankart Oslo / FKB',
        'MGS Oslo (kommunale data)',
    ]

    # Layout: to kolonner med tittel over hver
    col_w = Inches(5.2)
    gap = Inches(0.6)
    total = 2 * col_w + gap
    start_x = (SLIDE_W - total) / 2
    head_y = Inches(2.1)
    add_text(s, start_x, head_y,
             col_w, Inches(0.5),
             'Allerede på plass', size=16, bold=True,
             color=OSLO['morkegron'])
    add_text(s, start_x + col_w + gap, head_y,
             col_w, Inches(0.5),
             'Foreslåtte fremtidige kilder', size=16, bold=True,
             color=OSLO['morkgra'])

    row_h = Inches(0.55)
    row_gap = Inches(0.15)
    y0 = Inches(2.7)

    for i, label in enumerate(existing):
        y = y0 + i * (row_h + row_gap)
        outlined_box(s, start_x, y, col_w, row_h,
                     border_color=OSLO['morkegron'], border_w=1.2,
                     fill=OSLO['lysgron'])
        add_text(s, start_x + Inches(0.3), y + Inches(0.13),
                 col_w - Inches(0.6), Inches(0.35),
                 f'✓   {label}', size=13, bold=True,
                 color=OSLO['morkegron'])

    for i, label in enumerate(future):
        y = y0 + i * (row_h + row_gap)
        x = start_x + col_w + gap
        outlined_box(s, x, y, col_w, row_h,
                     border_color=OSLO['morkbeige'], border_w=1.2)
        add_text(s, x + Inches(0.3), y + Inches(0.13),
                 col_w - Inches(0.6), Inches(0.35),
                 f'○   {label}', size=13, color=OSLO['svart'])

    page_number(s, 11, TOTAL)


def s13_rag_intro(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Hva er RAG — og hvorfor trenger vi det?',
                 subtitle='RAG = Retrieval-Augmented Generation • å hente kunnskap før KI svarer')

    col_w = Inches(5.3)
    gap = Inches(0.5)
    total = 2 * col_w + gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.3)
    col_h = Inches(3.2)

    # Uten RAG
    outlined_box(s, start_x, y, col_w, col_h,
                 border_color=OSLO['morkbeige'], border_w=1.5)
    add_text(s, start_x + Inches(0.4), y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.5),
             'Uten RAG (dagens prototype)',
             size=16, bold=True, color=OSLO['morkebla'])
    add_bulleted(s, start_x + Inches(0.4), y + Inches(1.0),
                 col_w - Inches(0.8), col_h - Inches(1.2),
                 [
                     'Svarer fra det KI «husker» fra treningen',
                     'Ikke alltid oppdatert',
                     'Vi ser ikke hvor påstander kommer fra',
                 ], size=13, color=OSLO['svart'], after_pt=10)

    # Med RAG
    x2 = start_x + col_w + gap
    outlined_box(s, x2, y, col_w, col_h,
                 border_color=OSLO['morkegron'], border_w=1.5,
                 fill=OSLO['lysgron'])
    add_text(s, x2 + Inches(0.4), y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.5),
             'Med RAG (foreslått neste steg)',
             size=16, bold=True, color=OSLO['morkegron'])
    add_bulleted(s, x2 + Inches(0.4), y + Inches(1.0),
                 col_w - Inches(0.8), col_h - Inches(1.2),
                 [
                     'Henter fra Oslo kommunes egne dokumenter',
                     'Alltid oppdatert når dokumentbasen oppdateres',
                     'Hver påstand har sporbar kildehenvisning',
                 ], size=13, color=OSLO['morkegron'], after_pt=10)

    # Eksempel
    eks = outlined_box(s, MARGIN, Inches(6.0),
                       SLIDE_W - 2 * MARGIN, Inches(0.7),
                       border_color=OSLO['morkbeige'], border_w=1.0,
                       fill=OSLO['lysbeige'])
    fill_box_text(eks,
                  'Eksempel: «Hva sier Veileder Biotoptak om grønne tak?»  →  RAG slår opp svaret direkte i PDFen.',
                  size=13, italic=True, color=OSLO['svart'],
                  align=PP_ALIGN.CENTER)

    page_number(s, 12, TOTAL)


def s14_rag_architecture(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Slik vil RAG fungere for Naturportrett',
                 subtitle='Fem trinn fra Oslo kommunes dokumenter til kildebelagt portrett')

    steps = [
        ('Kunnskaps-\nbase', 'PDFer, veiledere,\nOslo-rapporter'),
        ('Indeksering', 'Deler tekst i biter\nfor søk'),
        ('Søkemotor', 'Finner relevante\nbiter ved spørring'),
        ('KI', 'Skriver svar med\nbitene som grunnlag'),
        ('Portrett\nm/ kilder', 'Sporbart per\npåstand'),
    ]
    box_w = Inches(2.05)
    box_h = Inches(2.4)
    gap = Inches(0.35)
    total = 5 * box_w + 4 * gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.3)
    for i, (label, body) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        # Siste trinn med subtil grønn fyll
        fill = OSLO['lysgron'] if i == 4 else None
        border = OSLO['morkegron'] if i == 4 else OSLO['morkebla']
        outlined_box(s, x, y, box_w, box_h,
                     border_color=border, border_w=1.5, fill=fill)

        add_text(s, x + Inches(0.15), y + Inches(0.35),
                 box_w - Inches(0.3), Inches(1.0),
                 label, size=14, bold=True,
                 color=OSLO['morkebla'] if fill is None else OSLO['morkegron'],
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(1.3),
                 box_w - Inches(0.3), Inches(1.0),
                 body, size=11,
                 color=OSLO['svart'],
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

        if i < 4:
            ax = x + box_w
            ay = y + box_h / 2
            add_arrow(s, ax + Inches(0.05), ay,
                      ax + gap - Inches(0.05), ay,
                      color=OSLO['morkebla'], width=2.0)

    # Verdiboks
    val_y = Inches(5.4)
    solid_box(s, MARGIN, val_y, SLIDE_W - 2 * MARGIN, Inches(1.0),
              fill=OSLO['morkebla'])
    add_text(s, MARGIN + Inches(0.3), val_y + Inches(0.15),
             SLIDE_W - 2 * MARGIN - Inches(0.6), Inches(0.4),
             'Verdien for Oslo kommune',
             size=15, bold=True, color=OSLO['hvit'])
    add_text(s, MARGIN + Inches(0.3), val_y + Inches(0.55),
             SLIDE_W - 2 * MARGIN - Inches(0.6), Inches(0.4),
             'Hver påstand kan lenkes tilbake til en konkret rapport eller veileder.',
             size=12, color=OSLO['lysebla_acc'])

    page_number(s, 13, TOTAL)


def s15_precision(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Mulige tiltak for mer presisjon',
                 subtitle='Tre retninger vi kan utforske')

    actions = [
        ('Mer presise områder på kartet',
         'I stedet for en enkel 500 m sirkel: bruk faktiske eiendomsgrenser, registrerte naturtyper og reguleringsformål for å avgjøre hva som er innenfor prosjektet.'),
        ('Direkte oppslag i offentlige databaser',
         'I tillegg til RAG (som henter fra dokumenter), kan systemet hente sanntidsfakta direkte fra Artsdatabanken og kommunens registre. F.eks. «hva er dagens rødlistestatus for spisslønn?» — svar fra databasen, ikke fra KIs hukommelse.'),
        ('Fagvurdering integrert i flyten',
         'Knapp for «godkjent av fagperson», redigeringsverktøy, og visuell flagging av usikre påstander.'),
    ]
    y = Inches(2.2)
    row_h = Inches(1.5)
    row_gap = Inches(0.2)
    for i, (head, body) in enumerate(actions):
        row_y = y + i * (row_h + row_gap)
        outlined_box(s, MARGIN, row_y, SLIDE_W - 2 * MARGIN,
                     row_h, border_color=OSLO['morkebla'],
                     border_w=1.2)
        # Nummer-sirkel
        num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  MARGIN + Inches(0.35), row_y + Inches(0.35),
                                  Inches(0.8), Inches(0.8))
        num.fill.solid()
        num.fill.fore_color.rgb = OSLO['morkebla']
        num.line.fill.background()
        fill_box_text(num, str(i + 1), size=24, bold=True,
                      color=OSLO['hvit'], align=PP_ALIGN.CENTER)
        add_text(s, MARGIN + Inches(1.4), row_y + Inches(0.2),
                 SLIDE_W - 2 * MARGIN - Inches(1.6), Inches(0.4),
                 head, size=17, bold=True, color=OSLO['morkebla'])
        add_text(s, MARGIN + Inches(1.4), row_y + Inches(0.65),
                 SLIDE_W - 2 * MARGIN - Inches(1.6), row_h - Inches(0.75),
                 body, size=13, color=OSLO['svart'], line_spacing=1.4)

    page_number(s, 14, TOTAL)


def s16_uke_questions(prs):
    s = prs.slides.add_slide(BLANK)
    title_header(s, 'Spørsmål vi bør drøfte med UKE',
                 subtitle='Når vi senere vurderer implementering i Oslo kommunes systemer')

    questions = [
        'Hvilken sky/infrastruktur passer best for denne typen tjenester?',
        'Hvordan integrere med eksisterende saksbehandlersystemer?',
        'Hvilken pålogging er aktuelt (saksbehandlere vs. eksterne brukere)?',
        'Hvordan håndtere drift, oppetid og support?',
        'Hva er sikkerhetskravene for KI-baserte verktøy i kommunen?',
        'Hvilke krav stilles til logging og sporbarhet?',
    ]

    add_bulleted(s, MARGIN, Inches(2.3),
                 SLIDE_W - 2 * MARGIN, Inches(4.5),
                 questions, size=18, color=OSLO['svart'],
                 line_spacing=1.5, after_pt=14)

    page_number(s, 15, TOTAL)


def s17_next_steps(prs):
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, OSLO['morkebla'])

    add_text(s, MARGIN, Inches(0.7), SLIDE_W - 2 * MARGIN, Inches(0.8),
             'Veien videre', size=38, bold=True, color=OSLO['hvit'])
    accent_stripe(s, MARGIN, Inches(1.55), Inches(2.0),
                   color=OSLO['gron_acc'])

    # Fase 1 — Juli 2026
    add_text(s, MARGIN, Inches(2.0),
             SLIDE_W - 2 * MARGIN, Inches(0.5),
             'Fase 1  •  Juli 2026', size=18, bold=True,
             color=OSLO['gron_acc'])
    add_text(s, MARGIN, Inches(2.5),
             SLIDE_W - 2 * MARGIN, Inches(0.4),
             'Jeg jobber videre med prototypen:', size=14,
             color=OSLO['lysebla_acc'])
    add_bulleted(s, MARGIN, Inches(2.95),
                 SLIDE_W - 2 * MARGIN, Inches(1.6),
                 [
                     'Bygge RAG-systemet — sporbarhet og bedre presisjon',
                     'Integrere flere datakilder',
                     'Forbedre format på produktet (PDF, sidelayout)',
                     'Diskutere om vi skal gjennomføre brukertesting',
                 ], size=14, color=OSLO['hvit'], after_pt=8)

    # Fase 2 — Etter sommeren
    add_text(s, MARGIN, Inches(5.0),
             SLIDE_W - 2 * MARGIN, Inches(0.5),
             'Fase 2  •  Etter sommeren', size=18, bold=True,
             color=OSLO['gron_acc'])
    add_text(s, MARGIN, Inches(5.5),
             SLIDE_W - 2 * MARGIN, Inches(0.4),
             'Møte med UKE for å drøfte:', size=14,
             color=OSLO['lysebla_acc'])
    add_bulleted(s, MARGIN, Inches(5.95),
                 SLIDE_W - 2 * MARGIN, Inches(0.9),
                 [
                     'Implementering i Oslo kommunes systemer',
                     'Alternativt: UKE bygger systemet, og denne prototypen tjener som skisse',
                 ], size=14, color=OSLO['hvit'], after_pt=8)

    page_number(s, 16, TOTAL, dark=True)


# ──────────────────────────────────────────────────────────────────────────────
# Build
# ──────────────────────────────────────────────────────────────────────────────

def build():
    prs = make_presentation()
    global BLANK
    BLANK = prs.slide_layouts[6]

    builders = [
        s01_title, s02_why, s03_products, s04_portrait_content,
        s05_user_flow, s06_demo_address, s07_demo_naturportrett,
        # s08_demo_detail droppet — to skjermbilder er nok
        s09_architecture, s10_data_sources,
        s11_ai_role, s12_more_sources, s13_rag_intro,
        s14_rag_architecture, s15_precision, s16_uke_questions,
        s17_next_steps,
    ]
    for builder in builders:
        builder(prs)

    prs.save(str(OUTPUT))
    print(f'OK — generert {OUTPUT} ({len(builders)} slides)')


if __name__ == '__main__':
    build()
